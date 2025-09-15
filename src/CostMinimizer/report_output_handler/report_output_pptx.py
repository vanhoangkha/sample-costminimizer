# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: Apache-2.0

__author__ = "Samuel Lepetre"
__license__ = "Apache-2.0"

from ..constants import __tooling_name__

from ..report_output_handler.report_output_handler import ReportOutputHandlerBase
from ..genai_providers.bedrock import Bedrock
from ..genai_providers.genai_providers import GenAIProviders
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Pt
import pandas as pd
import ast
import re
import json
# import RGBColor
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pathlib import Path
import os


class ReportOutputPptxHandlerBase(ReportOutputHandlerBase):
    DEFAULT_name_of_report_for_genai = 'ce_services'
    DEFAULT_name_of_genai_pptx_template = f'{__tooling_name__}_recommendations.pptx'
    
    def __init__(self, app, completed_reports, completion_time, determine_report_directory=True, create_directory_structure=False) -> None:
        super().__init__(app, completed_reports, completion_time, determine_report_directory, create_directory_structure)

        self.prs = None
        self.found_template_file = False
        self.default_chart_slide = 9
        self.default_number_format = '"$"#,##0'


    def create_directory_structure(self):
        self.make_report_directory_structure()
    
    def create_presentation(self):
        '''Load presentation template'''

        self.report_pptx_file = Path()
        try:
            # First attempt to find reports folder
            self.report_pptx_file = Path(os.getcwd()) / 'sample' / ReportOutputPptxHandlerBase.DEFAULT_name_of_genai_pptx_template  # Path to your existing template
            # check if file report_pptx_file exists
            if not os.path.exists(self.report_pptx_file):
                raise FileNotFoundError(f"File not found: {self.report_pptx_file}")
        except (OSError, FileNotFoundError):
            try:
                # Second attempt in src directory
                self.report_pptx_file = Path(os.getcwd()) / "src" / __tooling_name__ / 'sample' / ReportOutputPptxHandlerBase.DEFAULT_name_of_genai_pptx_template  # Path to your existing template
                # check if file report_pptx_file exists
                if not os.path.exists(self.report_pptx_file):
                    raise FileNotFoundError(f"File not found: {self.report_pptx_file}")
            except (OSError, FileNotFoundError) as e:
                self.logger.error(f'Unable to find the reports folder, either under {os.getcwd()} or src/')
                raise RuntimeError("Reports directory not found") from e

        self.found_template_file = True
        self.prs = Presentation(self.report_pptx_file)

    def get_data(self, report_name):
        '''obtain report data'''
        for report in self.completed_reports:
            if report.name() == report_name:
                return report.report_result
        return None
 
    def format_currency(self, value):
        '''format the values as currency'''
        return "${:,.0f}".format(value)

    def get_time_periods(self, data) -> list:
        '''return a list of of time_periods'''
        
        return data.keys().tolist()

    def prepare_chart_presentation(self, time_periods, number_format='"$"#,##0'):
        # Prepare the chart data
        chart_data = CategoryChartData(number_format=number_format)
        chart_data.categories = time_periods

        return chart_data
    
    def get_default_chart_slide(self, default_chart_slide_layout=3):
        '''return default slide along with default slide width and height '''
        slide_layout = self.prs.slide_layouts[default_chart_slide_layout]  # Using the second layout (0-based index)
        slide = self.prs.slides.add_slide(slide_layout)

        # Define slide dimensions (inches)
        chart_width = int(self.prs.slide_width)
        chart_height = int(self.prs.slide_height)

        return slide, chart_width, chart_height
    
    def center_chart_slide_dimensions(self, title_top, title_height, slide_width, slide_height, chart_width, chart_height):
        '''return demensions for centered chart on slide'''
        
        #obtain slide height and width
        slide_width = slide_width
        slide_height = slide_height

        min_y = title_top + title_height
        max_height = slide_height - min_y

        cx = chart_width
        if chart_height > max_height:
            cy = max_height
        else:
            cy = chart_height

        # Calculate dimensions to center the chart
        x = int((slide_width - cx) / 2)
        y = int((slide_height - cy) / 2)
        if y < min_y:
            y = min_y

        return x, y, cx, cy 
    
    def set_chart_title(self, chart, title):
        '''set chart title'''
        chart.has_title = True
        chart.chart_title.text_frame.text = title

        return chart
    
    def set_slide_title(self, slide, title):
        '''set slide title'''
        title_placeholder = slide.shapes.title
        if title_placeholder:
            # Adjust the slide title text
            title_placeholder.text = title
            # change title color to orange dark
            title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 153, 0)
    
    def set_default_legend(self, chart, pos=XL_LEGEND_POSITION.CORNER, font_size=8):
        '''set legend attributes'''
        chart.has_legend = True
        chart.legend.position = pos

        # Adjust chart to ensure the legend does not overlap
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(font_size)

        return chart

    def get_data_keys(self, data, data_key='account') -> list:
        '''return sorted list of a key '''
        return data.index.tolist()

    def add_series(self, data, top_ten_data, group_by_value='account') -> list:
        #data = self.accounts
        #top_ten_data = self.top_ten_accounts_data
        formatted_values = []
        for i in data:
            spend_values = top_ten_data.loc[i]
            
            # Format all values in the list
            formatted_values.append({i: [self.format_currency(value) for value in spend_values]})
            
            #update chart data series
            self.chart_data.add_series(f'{i}', spend_values)

        return formatted_values
     
    def top_ten(self, raw_data, data, group_by_value='account'):
        '''
        returns top 9 accounts (in total spend over time period) with spend
        grouped by time period.  The remaining accounts are aggregated into 
        "Others" 
        "'''
        
        return data['Data'].head(10)

    def last_months_spend_notes(self, data, record_type) -> str:
        '''return notes for last months spend'''
        last_month_spend = ""
        for entry in data:
            if entry['time_period'] == self.time_periods[-1]:
                last_month_spend += f"""{self.time_periods[-1]}: {entry[record_type]} : {self.format_currency(entry['amount'])}\n"""
            
        return last_month_spend

    def accounts_used_notes(self) -> str:
        '''return notes for accounts used'''
        accounts_used = ""
        for account in self.accounts:
                accounts_used += f"""{account}\n"""

        return accounts_used
    
    def generate_notes_text(self, note_frame, report_type, accounts, regions, excluded_charges,  last_month_spend, pre='', post=''):
                    
        if report_type == 'payer':
            report_type_text = f'This slide is based on all accounts under the payer: {accounts.split()[0]}'
        
        note_frame.text = f"""
Internal Notes:

{report_type_text}

Number of Accounts Used:
{len(accounts.split()[0])}

Selected Regions:
{regions}

Excluded Charge Types: 
{excluded_charges}

Last Months Spend:
{last_month_spend}
"""
        
        return note_frame

    def save_presentation(self, report_directory):
        '''save presentation'''
        self.make_report_directory_structure()
        file_name = report_directory / 'powerpoint_report/' / ReportOutputPptxHandlerBase.DEFAULT_name_of_genai_pptx_template

        # Construct the full output file path
        self.output_file = self.pptx_directory / file_name

        self.prs.save(self.output_file)


class ReportOutputPptxHandler(ReportOutputPptxHandlerBase):

    def __init__(self, appConfig, completed_reports, completion_time, determine_report_directory=True, create_directory_structure=False) -> None:
        super().__init__(appConfig, completed_reports, completion_time, determine_report_directory, create_directory_structure)

        self.selected_regions = self.appConfig.selected_regions

        if 'global' in self.selected_regions:
            #gloabl is stripped out from k2 queries currently
            self.selected_regions.remove('global')

        self.genai_provider = GenAIProviders()
        self.gen_ai_reccomendations_client = self.genai_provider.provider.client
        self.trend_spend_by_service = []
        self.analyzed_recommendations = []

    def compute_account_spend_summary(self):
        """
        Compute the total spend by account and the number of services generating spending.

        This function uses the 'ce_accounts' data to calculate the total spend for each account
        and count the number of unique services that generated spending for that account.

        Returns:
        list of dict: Each dictionary contains 'account_id', 'total_spend', and 'num_services'.
                      Returns None if no account data is available.
        """
        self.logger.info("Computing account spend summary")
        account_data = self.get_data('ce_accounts')
        if not account_data:
            self.logger.warning("No account data available")
            return [], [], 0, 0, 0, 0, {}

        list_accounts = account_data[0]['Data'].index.tolist()
        list_months = account_data[0]['Data'].keys().tolist()
        sum_last_month = account_data[0]['Data'].iloc[:,-1].sum()
        mean_last_month = account_data[0]['Data'].iloc[:,-1:].mean().sum()
        sum_last_6_months = account_data[0]['Data'].iloc[:,-6:].sum().sum()
        mean_sum_last_6_months = account_data[0]['Data'].iloc[:,-6:].sum().mean()

        account_summary = {}
        for account in list_accounts:
            account_summary[account] = dict()
            account_summary[account]['sum_spend'] = account_data[0]['Data'].loc[account].sum()
            account_summary[account]['mean_spend'] = account_data[0]['Data'].loc[account].mean()
            account_summary[account]['min_spend'] = account_data[0]['Data'].loc[account].mean()
            account_summary[account]['max_spend'] = account_data[0]['Data'].loc[account].mean()
            account_summary[account]['std_spend'] = account_data[0]['Data'].loc[account].std()

        self.logger.info(f"Computed spend summary for {len(account_summary)} accounts")
        return list_accounts, list_months, sum_last_month, mean_last_month, sum_last_6_months, mean_sum_last_6_months, account_summary

    def get_trend_spend_by_service_recommendations( self, gen_ai_client) -> list:
        
        msg=f"Getting trend by service AI recommendations: generate as many insights as you can about the monthly spend trends by service..."
        self.logger.info(msg)
        self.appConfig.console.print('\n'+msg)
        self.services_report_name = ReportOutputPptxHandlerBase.DEFAULT_name_of_report_for_genai
        self.services_data = self.get_data(self.services_report_name)

        # For each services_data, run the gen_ai_client to get a recommendation
        for service_data in self.services_data:
            df = pd.DataFrame.from_dict(service_data['Data'], orient='columns')
            if self.services_data is not None:
                gen_ai_data_list = gen_ai_client.execute( self.get_gen_ai_prompt('service_trends'), df, 'csv', True, 'dataframe')
                self.trend_spend_by_service.append(gen_ai_data_list)
            else:
                gen_ai_data_list = None
        return gen_ai_data_list

    def get_analyze_recommendations( self, gen_ai_client) -> list:
        msg=f'Getting service specific AI recommendations: create a list of talking points for the recommendations...'
        self.logger.info(msg)
        self.appConfig.console.print(msg)

        report_file = self.report_directory / self.appConfig.report_file_name
        gen_ai_data_list =  gen_ai_client.execute( gen_ai_client.get_gen_ai_prompt('recommendations'), report_file, 'xlsx', False, 'file')
        
        return gen_ai_data_list

    def create_disclaimer_slide(self):
        if self.found_template_file:
            slide, chart_width, chart_height = self.get_default_chart_slide()
            self.set_slide_title( slide, 'Disclaimer')

            # position and size of the bow should represent 90% of the slide size
            x, y, cx, cy = self.center_chart_slide_dimensions( slide.shapes.title.top, slide.shapes.title.height, chart_width, chart_height, int(chart_width*0.9), int(chart_height*0.7))

            text_box = slide.shapes.add_textbox(x, y, cx, cy)
            text_frame = text_box.text_frame
            # Enable word wrap and auto-fit text to shape
            text_frame.word_wrap = True
            text_frame.auto_fit = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # Add a text frame with bullet points
            bullet_points = [
                "- This presentation is to be considered beta.  Please double check values against Cost Explorer",
                "- We need your feedback.  What do you want to see in the presentation.",
                "- We will be adding in near term and long term recommendations as a roadmap item"  
            ]

            for bullet in bullet_points:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.level = 1  # This sets the bullet level, 0 being the top level


            # Add text to the notes section of the slide
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = ""
    
    def create_summary_slide(self):
        '''create summary slide'''
        if self.found_template_file:
            slide,  chart_width, chart_height = self.get_default_chart_slide()
            self.set_slide_title( slide, 'Cost Optimization Summary')

            # position and size of the bow should represent 90% of the new_slide size
            x, y, cx, cy = self.center_chart_slide_dimensions( slide.shapes.title.top, slide.shapes.title.height, chart_width, chart_height, int(chart_width*0.9), int(chart_height*0.7))

            text_box = slide.shapes.add_textbox(x, y, cx, cy)
            text_frame = text_box.text_frame
            # Enable word wrap and auto-fit text to shape
            text_frame.word_wrap = True
            text_frame.auto_fit = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            list_accounts, list_months, sum_last_month, mean_last_month, sum_last_6_months, mean_sum_last_6_months, results_per_accounts = self.compute_account_spend_summary()

            # Add a text frame with bullet points
            bullet_points = [
f"- AWS usage evaluated for cost optimization opportunities against {len(list_accounts)} accounts :",
f"         {list_accounts}",
f"- Regions for the following months : {list_months}",
f"- Spend over the last month : ${round(sum_last_month)}",
f"- Mean spend over the last month per account : ${round(mean_last_month)}",
f"- Spend over the last 6 months : ${round(sum_last_6_months)}",
f"- average spend over the last 6 months : ${round(mean_sum_last_6_months)}",
"- Note, these are best-case scenario estimates",
"- All recommendations given include a spreadsheet with corresponding resource data"
            ]

            for bullet in bullet_points:
                p = text_frame.add_paragraph()
                # Format as bullet point
                p.text = bullet
                p.font.size = Pt(20)
                #p.bullet.visible = True
                p.level = 1  # This sets the bullet level, 0 being the top level

            # Add text to the notes section of the slide
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = ""

    def create_spend_by_accounts_slide(self):
        '''create spend by accounts slide'''
        
        self.accounts_report_name = 'ce_accounts'
        self.accounts_data = self.get_data(self.accounts_report_name)  
        
        if self.accounts_data:
            if len(self.accounts_data) > 0 and isinstance(self.accounts_data, list):
                self.top_ten_accounts_data = self.top_ten(self.accounts_data, self.accounts_data[0], 'account_name')  
            
            # generate a list from the first row and of the dataframe
            self.time_periods = self.get_time_periods(self.top_ten_accounts_data)

            # generate a list from the first row and of the dataframe
            self.accounts = self.get_data_keys(self.top_ten_accounts_data, 'account_name')

            # generate a list from the first row and of the dataframe
            self.chart_data = self.prepare_chart_presentation(self.time_periods)

            formatted_values = self.add_series(self.accounts, self.top_ten_accounts_data, 'account_name')

            if self.found_template_file:
                slide, chart_width, chart_height = self.get_default_chart_slide()
                self.set_slide_title(slide, 'Trend Spend by Accounts')

                x, y, cx, cy = self.center_chart_slide_dimensions( slide.shapes.title.top, slide.shapes.title.height, chart_width, chart_height, int(chart_width*0.9), int(chart_height*0.7))

                chart = slide.shapes.add_chart( XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, self.chart_data).chart

                # set the title of the chart in yellow background
                chart = self.set_chart_title(chart, 'Spend by Accounts')

                # Change the background color of the chart title
                title = chart.chart_title
                fill = title.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
                # Change the foreground color of the chart title
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

                chart = self.set_default_legend(chart, XL_LEGEND_POSITION.CORNER, 8)

                # shift the chat so that it is the center of the screen
                chart.left = x
                chart.top = y

                # Add data labels to the chart
                for series in chart.series:
                    data_labels = series.data_labels
                    data_labels.number_format = '"$"#,##0'
                    series.data_labels.show_value = False
                    data_labels.font.size = Pt(6)  # Set the font size to 4 points
            
                # Add text to the notes section of the slide
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame

                if len(self.accounts_data[0]) > 0:
                    report_type_notes = 'payer'
                    excluded_charges_notes = 'excluded_charge_types'
                    regions_notes = f'Time periods {self.time_periods}'
                    last_month_spend_notes = f'payer: {self.accounts}'
                    accounts_used_notes = self.accounts_used_notes()

                self.generate_notes_text(notes_text_frame, report_type_notes, accounts_used_notes, regions_notes, excluded_charges_notes, last_month_spend_notes, pre='', post='')

    def create_spend_by_services_slide(self):
        
        self.services_report_name = 'ce_services'
        self.services_data = self.get_data(self.services_report_name)
        
        if self.services_data:
            if len(self.services_data) > 0 and isinstance(self.services_data, list):
                self.top_ten_services_data = self.top_ten(self.services_data, self.services_data[0], 'service') 
            
            # generate a list from the first row and of the dataframe
            self.time_periods = self.get_time_periods(self.top_ten_services_data)

            # generate a list from the first row and of the dataframe
            self.services = self.get_data_keys(self.top_ten_services_data, 'service')

            self.chart_data = self.prepare_chart_presentation(self.time_periods)

            #self.add_series_services()
            formatted_values = self.add_series(self.services, self.top_ten_services_data, 'service')

            if self.found_template_file:
                slide, chart_width, chart_height = self.get_default_chart_slide()
                self.set_slide_title(slide, 'Trend Spend by Services')

                x, y, cx, cy = self.center_chart_slide_dimensions( slide.shapes.title.top, slide.shapes.title.height, chart_width, chart_height, int(chart_width*0.9), int(chart_height*0.7))

                # add chart with a size that represent 70% of the slide size and in the middle for the location
                chart = slide.shapes.add_chart( XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, self.chart_data).chart

                # Change the background color of the chart title
                title = chart.chart_title
                fill = title.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
                # Change the foreground color of the chart title
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

                chart = self.set_default_legend(chart, XL_LEGEND_POSITION.CORNER, 8)

                # Add data labels to the chart
                for series in chart.series:
                    data_labels = series.data_labels
                    data_labels.number_format = '"$"#,##0'
                    series.data_labels.show_value = False
                    data_labels.font.size = Pt(6)  # Set the font size to 4 points

                # Add text to the notes section of the slide
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame

                if len(self.services_data[0]) > 0:
                    report_type_notes = 'payer'
                    excluded_charges_notes = 'excluded_charge_types'
                    regions_notes = 'regions_selected'
                    last_month_spend_notes = 'account_name'
                    accounts_used_notes = self.accounts_used_notes()

                self.generate_notes_text(notes_text_frame, report_type_notes, accounts_used_notes, regions_notes, excluded_charges_notes, last_month_spend_notes, pre='', post='')

    def create_section_seperator_slide(self, title):
        if self.found_template_file:
            slide, chart_width, chart_height = self.get_default_chart_slide(default_chart_slide_layout=7)
            self.set_slide_title(slide, title)

    def convert_reco_per_domain_to_string_formatted(self, reco_per_domain):
        """
        Convert reco_per_domain into a custom formatted string representation.
        
        Args:
            reco_per_domain: Dictionary of recommendations per domain
            
        Returns:
            str: Formatted string representation of the recommendations
        """
        try:
            output = []
            
            # If it's a dictionary
            if isinstance(reco_per_domain, dict):
                for domain, recommendations in reco_per_domain.items():
                    output.append(f"Domain: {domain}")
                    if isinstance(recommendations, list):
                        for reco in recommendations:
                            output.append(f"  - {reco}")
                    else:
                        output.append(f"  - {recommendations}")
                    output.append("")  # Empty line between domains
                    
            # If it's a list
            elif isinstance(reco_per_domain, list):
                for item in reco_per_domain:
                    output.append(str(item))
                    output.append("")  # Empty line between items
                    
            else:
                return str(reco_per_domain)
                
            return "\n".join(output)
            
        except Exception as e:
            print(f"Error converting recommendations to string: {e}")
            return str(reco_per_domain)

    def domain_based_recommendations_slide(self, domain):
        '''create summary slide'''
        l_bullet_points = []

        for report in self.completed_reports:
            if report.domain_name() == domain:
                l_sum = report.set_estimate_savings(sum=True)
                if l_sum >= 0:
                    formatted_amount = "${:,.0f}".format(l_sum)
                    l_text = f"- {report.common_name()}: {formatted_amount}"
                    l_bullet_points.append( l_text)

        if len(l_bullet_points) > 0:
            if self.found_template_file:
                slide, chart_width, chart_height = self.get_default_chart_slide()
                self.set_slide_title(slide, f'{domain}: Optimization Recommendations')

                # position and size of the bow should represent 90% of the new_slide size
                x, y, cx, cy = self.center_chart_slide_dimensions( slide.shapes.title.top, slide.shapes.title.height, chart_width, chart_height, int(chart_width*0.9), int(chart_height*0.7))

                text_box = slide.shapes.add_textbox(x, y, cx, cy)
                text_frame = text_box.text_frame
                # Enable word wrap and auto-fit text to shape
                text_frame.word_wrap = True
                text_frame.auto_fit = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                for bullet in l_bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(20)
                    p.level = 1  # This sets the bullet level, 0 being the top level

                # Add text to the notes section of the slide
                if self.appConfig.arguments_parsed.genai_recommendations:
                    reco_per_domain = self.add_genai_recommendations( self.trend_spend_by_service, domain)
                    notes_slide = slide.notes_slide
                    notes_text_frame = notes_slide.notes_text_frame
                    try:
                        notes_text_frame.text = self.convert_reco_per_domain_to_string_formatted(reco_per_domain)
                    except:
                        notes_text_frame.text = f'No recommendations for {domain}'


    def extract_dicts_from_string(self, s):
        # Remove any leading/trailing whitespace and ensure the string starts with '['
        s = s.strip()
        if not s.startswith('['):
            s = '[' + s

        # Ensure the string ends with '}]'
        if not s.endswith('}]'):
            s = s.rstrip(',') + '}]'

        # Replace single quotes with double quotes
        #s = s.replace("'", '"')

        # Use regex to find all dictionary-like structures
        dict_pattern = r'\{[^{}]*\}'
        dict_strings = re.findall(dict_pattern, s)

        dicts = []
        for dict_str in dict_strings:
            try:
                # Attempt to parse each dictionary string
                d = json.loads(dict_str)
                dicts.append(d)
            except json.JSONDecodeError as e:
                print(f"Error parsing dictionary: {e}")
                # Attempt to fix common issues
                fixed_dict_str = dict_str.replace('\n', '').replace('\\', '\\\\')
                try:
                    d = json.loads(fixed_dict_str)
                    dicts.append(d)
                except json.JSONDecodeError:
                    print(f"Failed to parse dictionary: {dict_str}")

        return dicts

    def add_genai_recommendations(self, recommendations, domain):
        bullet_points = recommendations
        result = []

        for bullet in bullet_points:
            
            # if domain is in bullet no matter the caps of the domain name
            # add it to the result list
            d = domain.lower()
            try:
                l = ast.literal_eval(bullet.replace('\n',''))['technical domain']
                b = l.lower()
                for item in l:
                    b = item['technical domain'].lower()
                    if (d in b):
                        result.append(item)
            except:
                try:
                    l = self.extract_dicts_from_string(bullet.replace('\n',''))
                    for item in l:
                        b = item['technical domain'].lower()
                        if (d in b):
                            result.append(item)
                    return result
                except:
                    return result
        return result

    def extract_dicts_from_string_regex(self, s):
        """
        Extract and parse dictionaries from a string using regex.
        This version can find dictionaries even if they're not in a proper list format.
        
        Args:
            s (str): Input string containing dictionaries
            
        Returns:
            list: List of dictionaries if found, empty list otherwise
        """
        try:
            # Use regex to find all dictionary-like structures
            dict_pattern = r'\{[^{}]*\}'
            dict_strings = re.findall(dict_pattern, s)
            
            dicts = []
            for dict_str in dict_strings:
                try:
                    # First try json.loads()
                    d = json.loads(dict_str)
                    dicts.append(d)
                except json.JSONDecodeError:
                    try:
                        # Try ast.literal_eval() if json.loads() fails
                        d = ast.literal_eval(dict_str)
                        dicts.append(d)
                    except (SyntaxError, ValueError):
                        # If both fail, try cleaning the string
                        dict_str = re.sub(r'\s+', ' ', dict_str).strip()
                        try:
                            d = ast.literal_eval(dict_str)
                            dicts.append(d)
                        except (SyntaxError, ValueError):
                            print(f"Failed to parse dictionary: {dict_str}")
                            
            return dicts
            
        except Exception as e:
            print(f"Error extracting dictionaries: {e}")
            return []

    def genai_recommendations_slide(self, recommendations):
        '''create slides for each technical domain with recommendations'''
        from collections import defaultdict

        # Group recommendations by technical domain
        domain_recommendations = defaultdict(list)

        # test if recommendations is not empty, else return
        if len(recommendations) == 0:
            return

        for item in recommendations:
            if type(item) == dict:
                l = item
            else:
                try:
                    l = ast.literal_eval(item)
                except:
                    try:
                        l = self.extract_dicts_from_string_regex(item)
                    except:
                        self.appConfig.console.print(f'Could not parse recommendation: {item}')
                        return

            domain_recommendations[l['technical domain']].append(l)

        for domain, items in domain_recommendations.items():

            for item in items:
                slide, chart_width, chart_height = self.get_default_chart_slide()
                self.set_slide_title(slide, f'GenAI: {domain} Recommendations')

                # position and size of the bow should represent 90% of the new_slide size
                x, y, cx, cy = self.center_chart_slide_dimensions( slide.shapes.title.top, slide.shapes.title.height, chart_width, chart_height, int(chart_width*0.9), int(chart_height*0.7))

                text_box = slide.shapes.add_textbox(x, y, cx, cy)
                text_frame = text_box.text_frame
                # Enable word wrap and auto-fit text to shape
                text_frame.word_wrap = True
                text_frame.auto_fit = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                p = text_frame.add_paragraph()
                p.text = f"Service: {item['service']}"
                p.font.size = Pt(24)
                p.font.bold = True
                p.level = 0

                p = text_frame.add_paragraph()
                p.text = item['recommendation data']
                p.font.size = Pt(20)
                p.level = 1

            # Add text to the notes section of the slide
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = f"Recommendations for {domain}"


    def fill_in_ppt_report(self, report_directory):

        try:
            self.create_disclaimer_slide()
        except:
            msg='ERROR : Unable to create disclaimer slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.create_summary_slide()
        except Exception as e:
            msg='ERROR : Unable to create summary slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.create_spend_by_accounts_slide()
        except:
            msg='ERROR : Unable to create spend by accounts slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)
        
        try:
            self.create_spend_by_services_slide()
        except:
            msg='ERROR : Unable to create spend by services slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)
        
        try:
            self.create_section_seperator_slide('TA, CO, CUR Recommendations')
        except:
            msg='ERROR : Unable to create separator slide for TA, CO, CUR Recommendations'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.domain_based_recommendations_slide('COMPUTE')
        except:
            msg='ERROR : Unable to create compute recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)
        
        try:
            self.domain_based_recommendations_slide('DATABASE')
        except:
            msg='ERROR : Unable to create database recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)
        
        try:
            self.domain_based_recommendations_slide('STORAGE')
        except:
            msg='ERROR : Unable to create storage recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.domain_based_recommendations_slide('NETWORK')
        except:
            msg='ERROR : Unable to create networking & content delivery recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.domain_based_recommendations_slide('ML')
        except:
            msg='ERROR : Unable to create machine learning recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.domain_based_recommendations_slide('MIGRATION_TRANSFER')
        except:
            msg='ERROR : Unable to create migration and transfer recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)
        
        try:
            self.domain_based_recommendations_slide('MANAGEMENT_GOVERNANCE')
        except:
            msg='ERROR : Unable to create management & governance recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.domain_based_recommendations_slide('ANALYTICS')
        except:
            msg='ERROR : Unable to create analytics recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.domain_based_recommendations_slide('APPLICATION_INTEGRATION')
        except:
            msg='ERROR : Unable to create application integration recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.create_section_seperator_slide('GenAI Recommendations')
        except:
            msg='ERROR : Unable to create separator slide for GenAI Recommendations'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.genai_recommendations_slide( self.analyzed_recommendations)
        except:
            msg='ERROR : Unable to create genai recommendations slide'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)

        try:
            self.save_presentation( report_directory)
            if self.appConfig.mode == 'cli':
                msg=f'\n[green]Powerpoint presentation saved to: [yellow]{self.output_file.resolve()}\n'
                self.appConfig.console.print(msg)
            msg=f'Powerpoint presentation saved to: {self.output_file.resolve()}\n'
            self.logger.info(msg)
        except:
            msg=f'ERROR : Unable to save presentation into : {report_directory}'
            self.logger.info(msg)
            self.appConfig.console.print('[red]'+msg)


    def export_to_csv(self, data, filename):
        # Implement CSV export logic here
        try:
            df = pd.DataFrame(data)
            df.to_csv(filename, index=False)
        finally:
            return

    def run(self):
        # test if genai-recommendations argument is true
        if self.appConfig.arguments_parsed.genai_recommendations:
            
            bedrock_gen_ai = Bedrock( self.appConfig)

            # test if '--ce' i.e. CostExplorer as argument is passed when the tooling is launched
            if self.appConfig.arguments_parsed.ce:

                self.trend_spend_by_service = self.get_trend_spend_by_service_recommendations(bedrock_gen_ai)

                # export self.trend_spend_by_service result into a csv file
                if len(self.trend_spend_by_service) > 0:
                    self.export_to_csv( self.trend_spend_by_service, self.report_directory / 'powerpoint_report/trend_spend_by_service_recommendations.csv')

            self.analyzed_recommendations = self.get_analyze_recommendations(bedrock_gen_ai)

            # export self.analyzed_recommendations result into a csv file
            if len(self.analyzed_recommendations) > 0:
                self.export_to_csv(self.analyzed_recommendations, self.report_directory / 'powerpoint_report/analyzed_recommendations.csv')

        # Create PPTX presentation file
        self.create_presentation()
